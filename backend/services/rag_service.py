import os
import re
from typing import List

def extract_text_from_pdf(filepath: str) -> str:
    """Extracts text from PDF file using pypdf."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        text = ""
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text += f"\n--- Page {page_num + 1} ---\n" + page_text
        if not text.strip():
            # If PDF has no extractable text layer (e.g. scanned image)
            text = f"Document: {os.path.basename(filepath)} (Total Pages: {len(reader.pages)})"
        return text
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return f"Document: {os.path.basename(filepath)}"

def extract_text_from_docx(filepath: str) -> str:
    """Extracts text from DOCX file."""
    try:
        import docx
        doc = docx.Document(filepath)
        paragraphs = [para.text for para in doc.paragraphs if para.text and para.text.strip()]
        # Also extract table text
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)
    except Exception as e:
        print(f"DOCX extraction notice: {e}, trying raw text extraction fallback...")
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                raw = f.read()
                # Extract printable ASCII/Unicode chunks
                printable = "".join([c for c in raw if c.isprintable() or c in "\n\r\t"])
                return printable[:15000]
        except Exception:
            return f"Document: {os.path.basename(filepath)}"

def extract_text_from_pptx(filepath: str) -> str:
    """Extracts text from PPTX slide deck."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text += f"\n--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_texts)
        if not text.strip():
            text = f"Presentation: {os.path.basename(filepath)} (Total Slides: {len(prs.slides)})"
        return text
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return f"Presentation: {os.path.basename(filepath)}"

def extract_document_content(filepath: str) -> str:
    """Extracts full readable text based on file format."""
    ext = os.path.splitext(filepath)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext in [".docx", ".doc"]:
        return extract_text_from_docx(filepath)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(filepath)
    elif ext in [".txt", ".md", ".markdown", ".csv", ".json", ".rtf", ".html", ".htm", ".log"]:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            print(f"Text read error for {filepath}: {e}")
            return f"Document: {os.path.basename(filepath)}"
    else:
        # Generic fallback for any other text-based or binary file
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                if content and any(c.isprintable() for c in content):
                    return content
        except Exception:
            pass
        return f"Educational Material: {os.path.basename(filepath)}"
